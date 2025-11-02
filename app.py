import streamlit as st
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer
import io
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import time
import re


st.set_page_config(page_title="Multi-Document Summarization", layout="wide")

st.title("🔤 Multi-Document Summarization")
st.write("Fine-tuned Pegasus Transformer Model")

# --- SUMMARY LENGTH CONTROLS ---
with st.sidebar:
    st.header("⚙️ Summary Settings")
    summary_length = st.slider(
        "Summary Length (words approx.)",
        min_value=50, max_value=500, value=150, step=50,
        help="Adjust the desired length of the summary"
    )
    
    processing_mode = st.radio(
        "Processing Mode:",
        ["🚀 Fast Mode", "🎯 Quality Mode"],
        help="Fast Mode: Quick combined summary | Quality Mode: Individual + Combined (slower)"
    )
    
    st.info(f"📝 Summary will be ~{summary_length} words")


# --- TEXT CLEANING FUNCTION ---
def clean_summary_text(text):
    """Clean and format summary text for better readability"""
    if not text:
        return ""
    
    # Remove <n> tokens
    text = re.sub(r'<n>', ' ', text)
    
    # Fix multiple periods/ellipses
    text = re.sub(r'\.\.+', '.', text)
    text = re.sub(r'\s*\.\s*\.\s*', '. ', text)
    
    # Fix spacing issues
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\s+([.,!?])', r'\1', text)
    
    # Fix sentence capitalization
    sentences = re.split(r'([.!?]+)', text)
    formatted_sentences = []
    for i in range(0, len(sentences)-1, 2):
        sentence = sentences[i].strip()
        if sentence:
            # Capitalize first letter
            sentence = sentence[0].upper() + sentence[1:] if len(sentence) > 1 else sentence.upper()
            punct = sentences[i+1] if i+1 < len(sentences) else '.'
            formatted_sentences.append(sentence + punct)
    
    # Join and clean final text
    text = ' '.join(formatted_sentences)
    text = re.sub(r'\s+', ' ', text).strip()
    
    return text


def format_summary_with_paragraphs(text, words_per_paragraph=50):
    """Break long summary into readable paragraphs"""
    if not text:
        return ""
    
    # Split into sentences
    sentences = re.split(r'(?<=[.!?])\s+', text)
    
    paragraphs = []
    current_para = []
    word_count = 0
    
    for sentence in sentences:
        sentence_words = len(sentence.split())
        
        if word_count + sentence_words > words_per_paragraph and current_para:
            paragraphs.append(' '.join(current_para))
            current_para = [sentence]
            word_count = sentence_words
        else:
            current_para.append(sentence)
            word_count += sentence_words
    
    # Add remaining sentences
    if current_para:
        paragraphs.append(' '.join(current_para))
    
    return '\n\n'.join(paragraphs)


# --- TEXT EXTRACTION FUNCTIONS ---
def extract_text_from_pdf(file):
    """Extracts text content from a PDF file."""
    text = ""
    try:
        file_bytes = io.BytesIO(file.read())
        reader = PdfReader(file_bytes)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.error(f"Error reading PDF: {e}")
    return text.strip()


def extract_text_from_pptx(file):
    """Extracts text content from a PPTX file."""
    text = ""
    try:
        file_bytes = io.BytesIO(file.read())
        prs = Presentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f"Error reading PPTX: {e}")
    return text.strip()


def extract_text_from_docx(file):
    """Extracts text content from a DOCX file."""
    text = ""
    try:
        file_bytes = io.BytesIO(file.read())
        doc = Document(file_bytes)
        for paragraph in doc.paragraphs:
            if paragraph.text:
                text += paragraph.text + "\n"
    except Exception as e:
        st.error(f"Error reading DOCX: {e}")
    return text.strip()


# --- MODEL LOADING ---
@st.cache_resource(show_spinner=False)
def load_model():
    """Load model from Hugging Face Hub with caching"""
    try:
        with st.spinner("🔄 Loading model... This may take a minute"):
            model = AutoModelForSeq2SeqLM.from_pretrained("google/pegasus-cnn_dailymail")
            tokenizer = AutoTokenizer.from_pretrained("google/pegasus-cnn_dailymail")
        return model, tokenizer
    except Exception as e:
        st.error(f"Error loading model: {e}")
        return None, None


def summarize_text(text, model, tokenizer, max_length=150, min_length=None):
    """Generate summary for a single text with adjustable length and better quality"""
    if not text.strip():
        return "No content to summarize."
    
    try:
        # Calculate token lengths
        max_tokens = min(int(max_length * 1.5), 512)
        min_tokens = min_length if min_length else max(30, max_tokens // 3)
        
        # Tokenize input
        inputs = tokenizer(
            text, 
            return_tensors="pt", 
            max_length=1024, 
            truncation=True,
            padding=True
        )
        
        # Generate summary with optimized parameters
        summary_ids = model.generate(
            inputs["input_ids"],
            max_length=max_tokens,
            min_length=min_tokens,
            num_beams=4,
            length_penalty=1.0,
            early_stopping=True,
            no_repeat_ngram_size=3,
            repetition_penalty=1.2,
            temperature=1.0
        )
        
        # Decode and clean
        summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
        summary = clean_summary_text(summary)
        
        return summary
        
    except Exception as e:
        return f"Error generating summary: {str(e)}"


# --- OPTIMIZED SUMMARIZATION STRATEGIES ---
def fast_combined_summary(file_contents, model, tokenizer, max_length=150):
    """
    FAST MODE: Extract key sentences from each document and combine them
    """
    if not file_contents:
        return "No content to summarize.", []
    
    important_sentences = []
    
    for content in file_contents:
        if content.strip():
            sentences = re.split(r'(?<=[.!?])\s+', content)
            # Take first 3 and last 2 sentences from each document
            important_sentences.extend(sentences[:3])
            if len(sentences) > 5:
                important_sentences.extend(sentences[-2:])
    
    combined_text = " ".join([s.strip() for s in important_sentences if s.strip()])
    
    if combined_text:
        final_summary = summarize_text(combined_text, model, tokenizer, max_length)
        final_summary = format_summary_with_paragraphs(final_summary)
        return final_summary, important_sentences
    else:
        return "No meaningful content found.", []


def quality_combined_summary(file_contents, model, tokenizer, max_length=150):
    """
    QUALITY MODE: Proper two-stage summarization for better quality
    """
    if not file_contents:
        return "No content to summarize.", []
    
    individual_summaries = []
    
    # Stage 1: Summarize each document individually
    per_doc_length = max(80, max_length // len(file_contents)) if len(file_contents) > 1 else max_length
    
    for i, content in enumerate(file_contents):
        if content.strip():
            individual_summary = summarize_text(
                content, 
                model, 
                tokenizer, 
                max_length=per_doc_length,
                min_length=per_doc_length // 2
            )
            individual_summaries.append(individual_summary)
    
    if individual_summaries:
        # Stage 2: Combine individual summaries into final summary
        combined_text = " ".join(individual_summaries)
        
        final_summary = summarize_text(
            combined_text, 
            model, 
            tokenizer, 
            max_length=max_length,
            min_length=max(max_length // 2, 50)
        )
        
        # Format with paragraphs
        final_summary = format_summary_with_paragraphs(final_summary)
        
        return final_summary, individual_summaries
    else:
        return "No meaningful content found.", []


# --- SIMPLIFIED FILE PROCESSING ---
def process_single_file(file):
    """Extract content from a single file"""
    try:
        file.seek(0)
        
        if file.name.lower().endswith('.pdf'):
            content = extract_text_from_pdf(file)
        elif file.name.lower().endswith('.pptx'):
            content = extract_text_from_pptx(file)
        elif file.name.lower().endswith('.docx'):
            content = extract_text_from_docx(file)
        else:
            content = file.read().decode("utf-8")
        
        return content.strip() if content else ""
    except Exception as e:
        st.warning(f"❌ Error processing {file.name}: {e}")
        return ""


# Load model
model, tokenizer = load_model()

# --- STREAMLIT UI ---
if 'text_documents' not in st.session_state:
    st.session_state.text_documents = [""]
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = []


tab1, tab2 = st.tabs(["📝 Multiple Text Inputs", "📄 File Upload"])


with tab1:
    st.write("**Enter multiple documents below (one per box):**")
    
    col1, col2 = st.columns([1, 5])
    with col1:
        if st.button("➕ Add Document", key="add_doc_text"):
            st.session_state.text_documents.append("")
    with col2:
        if st.button("➖ Remove Document", key="remove_doc_text") and len(st.session_state.text_documents) > 1:
            st.session_state.text_documents.pop()
    
    for i in range(len(st.session_state.text_documents)):
        st.session_state.text_documents[i] = st.text_area(
            f"Document {i+1}:", 
            height=150, 
            key=f"doc_input_{i}",
            placeholder=f"Paste document {i+1} here...",
            value=st.session_state.text_documents[i]
        )

    summary_type_text = st.radio(
        "Summary Scope:",
        ["Current Tab (Text Inputs Only)", "Combined Tab Data (Text Inputs + Uploaded Files)"],
        horizontal=True,
        key="summary_type_text"
    )
    
    if st.button("🚀 Generate Summary", key="text_summary_btn"):
        with st.spinner("Generating summary..."):
            start_time = time.time()
            try:
                text_contents = [doc for doc in st.session_state.text_documents if doc.strip()]
                
                file_contents = []
                if summary_type_text == "Combined Tab Data (Text Inputs + Uploaded Files)" and st.session_state.uploaded_files:
                    for file in st.session_state.uploaded_files:
                        content = process_single_file(file)
                        if content:
                            file_contents.append(content)
                
                all_contents = text_contents + file_contents
                
                if not all_contents:
                    st.warning("Please provide some text or files to summarize.")
                else:
                    if processing_mode == "🚀 Fast Mode":
                        final_summary, intermediate_data = fast_combined_summary(
                            all_contents, model, tokenizer, summary_length
                        )
                        intermediate_label = "Key Sentences Used"
                    else:
                        final_summary, intermediate_data = quality_combined_summary(
                            all_contents, model, tokenizer, summary_length
                        )
                        intermediate_label = "Intermediate Summaries"
                    
                    end_time = time.time()
                    st.success(f"✅ Combined Summary Generated! (Time: {end_time - start_time:.2f}s)")
                    st.markdown("### 📋 Combined Summary")
                    
                    word_count = len(final_summary.split())
                    st.caption(f"Summary length: {word_count} words (target: {summary_length}) | Sources: {len(all_contents)} documents | Mode: {processing_mode}")
                    
                    # FIXED: Use st.info() to display summary (your original method)
                    st.info(final_summary)
                    
                    if intermediate_data:
                        with st.expander(f"🔍 {intermediate_label}"):
                            for i, data in enumerate(intermediate_data):
                                st.write(f"**Document {i+1}:**")
                                cleaned_data = clean_summary_text(str(data)) if isinstance(data, str) else str(data)
                                st.text(cleaned_data[:500] + "..." if len(cleaned_data) > 500 else cleaned_data)
                                st.divider()
                        
            except Exception as e:
                st.error(f"Error generating summary: {e}")


with tab2:
    st.write("**Upload files (.txt, .pdf, .pptx, .docx) for summarization:**")
    
    uploaded_files = st.file_uploader(
        "Choose files", 
        type=['txt', 'pdf', 'pptx', 'docx'], 
        accept_multiple_files=True,
        key="file_uploader"
    )

    if uploaded_files is not None:
        st.session_state.uploaded_files = uploaded_files
    
    if st.session_state.uploaded_files:
        st.write(f"**📁 {len(st.session_state.uploaded_files)} file(s) uploaded**")
        
        file_contents_data = []
        for file in st.session_state.uploaded_files:
            content = process_single_file(file)
            if content:
                file_contents_data.append({
                    'name': file.name,
                    'content': content,
                    'word_count': len(content.split())
                })
        
        with st.expander("📊 File Contents Preview"):
            for file_data in file_contents_data:
                st.write(f"**{file_data['name']}** - {file_data['word_count']} words")
                preview = file_data['content'][:500] + "..." if len(file_data['content']) > 500 else file_data['content']
                st.text(preview)
                st.divider()
        
        file_summary_type = st.radio(
            "Summary Scope:",
            ["Current Tab (Uploaded Files Only)", "Combined Tab Data (Uploaded Files + Text Inputs)"],
            horizontal=True,
            key="file_summary_type"
        )
        
        individual_summary_check = st.checkbox(
            "Also generate individual summaries for uploaded files?", 
            key="individual_check"
        )

        if st.button("🚀 Generate Summary", key="file_summary_btn"):
            with st.spinner("Processing files and generating summary..."):
                start_time = time.time()
                try:
                    file_contents_list = [file_data['content'] for file_data in file_contents_data]
                    
                    text_contents = []
                    if file_summary_type == "Combined Tab Data (Uploaded Files + Text Inputs)":
                        text_contents = [doc for doc in st.session_state.text_documents if doc.strip()]
                    
                    all_contents = file_contents_list + text_contents

                    if not all_contents:
                        st.warning("No readable content was found.")
                    else:
                        if processing_mode == "🚀 Fast Mode":
                            final_summary, intermediate_data = fast_combined_summary(
                                all_contents, model, tokenizer, summary_length
                            )
                            intermediate_label = "Key Sentences Used"
                        else:
                            final_summary, intermediate_data = quality_combined_summary(
                                all_contents, model, tokenizer, summary_length
                            )
                            intermediate_label = "Intermediate Summaries"
                        
                        end_time = time.time()
                        st.success(f"✅ Combined Summary Generated! (Time: {end_time - start_time:.2f}s)")
                        st.markdown("### 📋 Combined Summary")
                        
                        word_count = len(final_summary.split())
                        st.caption(f"Summary length: {word_count} words (target: {summary_length}) | Sources: {len(all_contents)} documents | Mode: {processing_mode}")
                        
                        # FIXED: Use st.info() to display summary properly
                        st.info(final_summary)
                        
                        if intermediate_data:
                            with st.expander(f"🔍 {intermediate_label}"):
                                for i, data in enumerate(intermediate_data):
                                    st.write(f"**Source {i+1}:**")
                                    cleaned_data = clean_summary_text(str(data)) if isinstance(data, str) else str(data)
                                    st.text(cleaned_data[:500] + "..." if len(cleaned_data) > 500 else cleaned_data)
                                    st.divider()
                        
                        # Individual summaries
                        if individual_summary_check:
                            st.divider()
                            st.markdown("### 📄 Individual File Summaries")
                            
                            for file_data in file_contents_data:
                                with st.spinner(f"Summarizing {file_data['name']}..."):
                                    individual_start = time.time()
                                    individual_summary = summarize_text(
                                        file_data['content'], 
                                        model, 
                                        tokenizer, 
                                        max_length=summary_length
                                    )
                                    individual_summary = format_summary_with_paragraphs(individual_summary)
                                    individual_end = time.time()
                                
                                st.markdown(f"**{file_data['name']}** ({file_data['word_count']} words)")
                                st.caption(f"Generated in {individual_end - individual_start:.2f}s | {len(individual_summary.split())} words")
                                # Use st.info() for display
                                st.info(individual_summary)
                                st.divider()
                                
                except Exception as e:
                    st.error(f"Error generating summary: {e}")


# --- MODEL METRICS ---
st.divider()
st.subheader("📊 Model Performance Metrics")

metrics = {
    "Pegasus (CNN/DailyMail)": {"ROUGE-1": 47.65, "ROUGE-2": 18.75, "ROUGE-L": 24.95},
    "TextRank": {"ROUGE-1": 43.83, "ROUGE-2": 7.97, "ROUGE-L": 34.13},
    "LSTM-Seq2Seq": {"ROUGE-1": 38.0, "ROUGE-2": 17.0, "ROUGE-L": 32.0}
}

st.json(metrics)
